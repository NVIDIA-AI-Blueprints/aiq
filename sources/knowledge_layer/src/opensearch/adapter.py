# SPDX-FileCopyrightText: Copyright (c) 2025-2026, NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
# http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""
OpenSearch adapter for the Knowledge Layer.

This backend stores one OpenSearch vector index per AIQ collection/session and
supports three authentication modes:
- none: self-hosted development clusters without authentication
- basic: self-hosted clusters with username/password
- sigv4: Amazon OpenSearch Service and Amazon OpenSearch Serverless
"""

from __future__ import annotations

import logging
import os
import re
import threading
import time
import uuid
from datetime import UTC
from datetime import datetime
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

from aiq_agent.knowledge.base import BaseIngestor
from aiq_agent.knowledge.base import BaseRetriever
from aiq_agent.knowledge.base import TTLCleanupMixin
from aiq_agent.knowledge.factory import register_ingestor
from aiq_agent.knowledge.factory import register_retriever
from aiq_agent.knowledge.schema import Chunk
from aiq_agent.knowledge.schema import CollectionInfo
from aiq_agent.knowledge.schema import ContentType
from aiq_agent.knowledge.schema import FileInfo
from aiq_agent.knowledge.schema import FileProgress
from aiq_agent.knowledge.schema import FileStatus
from aiq_agent.knowledge.schema import IngestionJobStatus
from aiq_agent.knowledge.schema import JobState
from aiq_agent.knowledge.schema import RetrievalResult

logger = logging.getLogger(__name__)

# @environment_variable OPENSEARCH_URL
# @category Knowledge Layer
# @type str
# @default http://localhost:9200
# @required false
# Base URL for self-hosted OpenSearch, Amazon OpenSearch, or OpenSearch Serverless.
DEFAULT_ENDPOINT = os.environ.get("OPENSEARCH_URL", "http://localhost:9200")

# @environment_variable OPENSEARCH_AUTH_TYPE
# @category Knowledge Layer
# @type str
# @default none
# @required false
# Auth mode for OpenSearch: none, basic, or sigv4.
DEFAULT_AUTH_TYPE = os.environ.get("OPENSEARCH_AUTH_TYPE", "none")

DEFAULT_INDEX_PREFIX = os.environ.get("OPENSEARCH_INDEX_PREFIX", "aiq")
DEFAULT_AWS_REGION = os.environ.get("AWS_REGION") or os.environ.get("AWS_DEFAULT_REGION", "us-east-1")
DEFAULT_AWS_SERVICE = os.environ.get("OPENSEARCH_AWS_SERVICE", "aoss")
DEFAULT_EMBED_MODEL = os.environ.get("AIQ_EMBED_MODEL", "nvidia/llama-nemotron-embed-vl-1b-v2")
DEFAULT_EMBED_BASE_URL = os.environ.get("AIQ_EMBED_BASE_URL", "https://integrate.api.nvidia.com/v1")
DEFAULT_VECTOR_FIELD = os.environ.get("OPENSEARCH_VECTOR_FIELD", "embedding")
DEFAULT_TEXT_FIELD = os.environ.get("OPENSEARCH_TEXT_FIELD", "content")
DEFAULT_EMBEDDING_DIM = int(os.environ.get("OPENSEARCH_EMBEDDING_DIM", "2048"))
DEFAULT_TIMEOUT = int(os.environ.get("OPENSEARCH_TIMEOUT", "120"))
DEFAULT_CHUNK_SIZE = int(os.environ.get("OPENSEARCH_CHUNK_SIZE", "1024"))
DEFAULT_CHUNK_OVERLAP = int(os.environ.get("OPENSEARCH_CHUNK_OVERLAP", "128"))
DEFAULT_INGESTION_MODE = os.environ.get("OPENSEARCH_INGESTION_MODE", "local")
DEFAULT_DASK_SCHEDULER_ADDRESS = os.environ.get("OPENSEARCH_DASK_SCHEDULER_ADDRESS") or os.environ.get(
    "NAT_DASK_SCHEDULER_ADDRESS"
)
DEFAULT_DASK_FILE_TRANSFER = os.environ.get("OPENSEARCH_DASK_FILE_TRANSFER", "bytes")
DEFAULT_AOSS_DELETE_MAX_BATCHES = int(os.environ.get("OPENSEARCH_AOSS_DELETE_MAX_BATCHES", "100"))
DEFAULT_AOSS_DELETE_BACKOFF_SECONDS = float(os.environ.get("OPENSEARCH_AOSS_DELETE_BACKOFF_SECONDS", "0.25"))

# Collection TTL settings, aligned with the other knowledge backends.
COLLECTION_TTL_HOURS = float(os.environ.get("AIQ_COLLECTION_TTL_HOURS", "24"))
TTL_CLEANUP_INTERVAL_SECONDS = int(os.environ.get("AIQ_TTL_CLEANUP_INTERVAL_SECONDS", "3600"))

SUMMARY_MAX_INPUT_CHARS = 4000
DEFAULT_BULK_BATCH_SIZE = 100
DEFAULT_EMBEDDING_BATCH_SIZE = 16
SUPPORTED_TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".log"}


def _utc_now() -> datetime:
    return datetime.now(tz=UTC)


def _sanitize_index_part(value: str, fallback: str = "default") -> str:
    """Convert a collection/prefix value into an OpenSearch-safe index name part."""
    normalized = re.sub(r"[^a-z0-9._-]+", "-", value.lower()).strip(".-_")
    if not normalized:
        normalized = fallback
    if normalized[0] in ("-", "_", "+"):
        normalized = f"x-{normalized.lstrip('-_+')}"
    return normalized


def _trim_index_name(index_name: str) -> str:
    """Ensure the physical index name stays within OpenSearch's length limit."""
    if len(index_name) <= 255:
        return index_name
    suffix = uuid.uuid5(uuid.NAMESPACE_URL, index_name).hex[:12]
    return f"{index_name[:242]}-{suffix}"


def _normalize_endpoint(endpoint: str, default_scheme: str) -> str:
    endpoint = str(endpoint).strip()
    if "://" not in endpoint:
        endpoint = f"{default_scheme}://{endpoint}"
    return endpoint.rstrip("/")


def _score_to_similarity(score: Any) -> float:
    """Normalize backend scores into the universal [0, 1] score contract."""
    try:
        value = float(score)
    except (TypeError, ValueError):
        return 0.0
    return max(0.0, min(1.0, value))


def _parse_timestamp(value: Any) -> datetime | None:
    if not value:
        return None
    if isinstance(value, datetime):
        return value
    if isinstance(value, str):
        try:
            return datetime.fromisoformat(value.replace("Z", "+00:00"))
        except ValueError:
            return None
    return None


def _generate_document_summary(text_content: str, file_name: str, llm=None) -> str | None:
    """Generate a one-sentence document summary using the configured LangChain LLM."""
    if llm is None or not text_content.strip():
        return None

    prompt = (
        "Summarize this uploaded document in one concise sentence for a research assistant. "
        "Focus on the document's topic and likely usefulness.\n\n"
        f"Document: {file_name}\n\n"
        f"Content excerpt:\n{text_content[:SUMMARY_MAX_INPUT_CHARS]}"
    )

    try:
        response = llm.invoke(prompt)
        summary = getattr(response, "content", response)
        summary_text = str(summary).strip()
        return summary_text[:500] if summary_text else None
    except Exception as e:
        logger.warning("Summary generation failed for %s: %s", file_name, e)
        return None


def _read_text_file(file_path: Path) -> list[tuple[str, int | None, dict[str, Any]]]:
    content = file_path.read_text(encoding="utf-8", errors="ignore")
    return [(content, None, {"file_type": file_path.suffix.lower().lstrip(".") or "text"})]


def _read_pdf_file(file_path: Path) -> list[tuple[str, int | None, dict[str, Any]]]:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError("PDF ingestion for OpenSearch requires pypdf. Install knowledge-layer[opensearch].") from e

    reader = PdfReader(str(file_path))
    pages: list[tuple[str, int | None, dict[str, Any]]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((text, idx, {"file_type": "pdf"}))
    return pages


def _read_docx_file(file_path: Path) -> list[tuple[str, int | None, dict[str, Any]]]:
    try:
        import docx2txt
    except ImportError as e:
        raise RuntimeError(
            "DOCX ingestion for OpenSearch requires docx2txt. Install knowledge-layer[opensearch]."
        ) from e

    return [(docx2txt.process(str(file_path)) or "", None, {"file_type": "docx"})]


def _read_pptx_file(file_path: Path) -> list[tuple[str, int | None, dict[str, Any]]]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError(
            "PPTX ingestion for OpenSearch requires python-pptx. Install knowledge-layer[opensearch]."
        ) from e

    presentation = Presentation(str(file_path))
    slides: list[tuple[str, int | None, dict[str, Any]]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        content = "\n".join(texts)
        if content.strip():
            slides.append((content, idx, {"file_type": "pptx", "slide_number": idx}))
    return slides


def _read_file_segments(file_path: str) -> list[tuple[str, int | None, dict[str, Any]]]:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix in SUPPORTED_TEXT_EXTENSIONS or not suffix:
        return _read_text_file(path)
    if suffix == ".pdf":
        return _read_pdf_file(path)
    if suffix == ".docx":
        return _read_docx_file(path)
    if suffix == ".pptx":
        return _read_pptx_file(path)
    raise RuntimeError(f"Unsupported file type for OpenSearch ingestion: {suffix}")


def _chunk_text(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """Chunk text using a word-count approximation that avoids tokenizer dependencies."""
    words = text.split()
    if not words:
        return []

    chunk_size = max(1, chunk_size)
    chunk_overlap = max(0, min(chunk_overlap, chunk_size - 1))
    step = chunk_size - chunk_overlap
    chunks = []

    for start in range(0, len(words), step):
        chunk_words = words[start : start + chunk_size]
        if not chunk_words:
            break
        chunks.append(" ".join(chunk_words))
        if start + chunk_size >= len(words):
            break

    return chunks


def _resolve_embedding_api_key(embed_base_url: str) -> str:
    api_key = os.environ.get("NVIDIA_API_KEY", "")
    is_hosted_nvidia = "integrate.api.nvidia.com" in (embed_base_url or "")
    if is_hosted_nvidia and not api_key:
        raise RuntimeError(
            "NVIDIA_API_KEY is required for the hosted NVIDIA embeddings API "
            "(embed_base_url contains integrate.api.nvidia.com). Either set "
            "NVIDIA_API_KEY or override AIQ_EMBED_BASE_URL to a self-hosted NIM endpoint."
        )
    return api_key


class _OpenSearchConfigMixin:
    """Shared OpenSearch configuration and client helpers."""

    config: dict[str, Any]
    _client: Any
    _client_lock: threading.RLock

    def _configure_opensearch(self) -> None:
        self.auth_type = str(self.config.get("auth_type", DEFAULT_AUTH_TYPE)).lower()
        raw_endpoint = self.config.get("endpoint") or self.config.get("opensearch_url") or DEFAULT_ENDPOINT
        self.endpoint = _normalize_endpoint(raw_endpoint, "https" if self.auth_type == "sigv4" else "http")
        self.username = self.config.get("username") or os.environ.get("OPENSEARCH_USERNAME")
        self.password = self.config.get("password") or os.environ.get("OPENSEARCH_PASSWORD")
        self.aws_region = self.config.get("aws_region", DEFAULT_AWS_REGION)
        self.aws_service = self.config.get("aws_service", DEFAULT_AWS_SERVICE)
        self.verify_certs = self.config.get("verify_certs", True)
        self.ca_certs = self.config.get("ca_certs") or os.environ.get("OPENSEARCH_CA_CERTS")
        self.timeout = self.config.get("timeout", DEFAULT_TIMEOUT)
        self.max_retries = self.config.get("max_retries", 3)
        self.retry_on_timeout = self.config.get("retry_on_timeout", True)
        self.index_prefix = _sanitize_index_part(self.config.get("index_prefix", DEFAULT_INDEX_PREFIX), "aiq")
        self.vector_field = self.config.get("vector_field", DEFAULT_VECTOR_FIELD)
        self.text_field = self.config.get("text_field", DEFAULT_TEXT_FIELD)
        self.embedding_dim = int(self.config.get("embedding_dim", DEFAULT_EMBEDDING_DIM))
        self.engine = self.config.get("engine", "faiss")
        self.space_type = self.config.get("space_type", "cosinesimil")
        self.m = int(self.config.get("m", 16))
        self.ef_construction = int(self.config.get("ef_construction", 512))
        self.ef_search = int(self.config.get("ef_search", 512))
        self.bulk_batch_size = int(self.config.get("bulk_batch_size", DEFAULT_BULK_BATCH_SIZE))
        self.aoss_delete_max_batches = int(self.config.get("aoss_delete_max_batches", DEFAULT_AOSS_DELETE_MAX_BATCHES))
        self.aoss_delete_backoff_seconds = float(
            self.config.get("aoss_delete_backoff_seconds", DEFAULT_AOSS_DELETE_BACKOFF_SECONDS)
        )
        self.allow_document_ids = self.config.get("allow_document_ids")
        if self.allow_document_ids is None:
            self.allow_document_ids = not (self.auth_type == "sigv4" and self.aws_service == "aoss")
        self.bulk_refresh = self.config.get("bulk_refresh")
        if self.bulk_refresh is None:
            self.bulk_refresh = False if self.auth_type == "sigv4" and self.aws_service == "aoss" else True

        self._client = None
        self._client_lock = threading.RLock()

    def _index_name_for_collection(self, collection_name: str) -> str:
        collection_part = _sanitize_index_part(collection_name, "default")
        return _trim_index_name(f"{self.index_prefix}-{collection_part}")

    def _create_client(self):
        try:
            from opensearchpy import OpenSearch
            from opensearchpy import RequestsHttpConnection
        except ImportError as e:
            raise RuntimeError(
                "OpenSearch dependencies not installed. Install with: knowledge-layer[opensearch]"
            ) from e

        parsed = urlparse(self.endpoint)
        use_ssl = parsed.scheme == "https"
        client_kwargs: dict[str, Any] = {
            "use_ssl": use_ssl,
            "verify_certs": self.verify_certs,
            "timeout": self.timeout,
            "max_retries": self.max_retries,
            "retry_on_timeout": self.retry_on_timeout,
            "connection_class": RequestsHttpConnection,
        }
        if self.ca_certs:
            client_kwargs["ca_certs"] = self.ca_certs

        if self.auth_type == "basic":
            if not self.username or not self.password:
                raise RuntimeError("OpenSearch basic auth requires username and password")
            client_kwargs["http_auth"] = (self.username, self.password)
        elif self.auth_type == "sigv4":
            try:
                import boto3
                from opensearchpy import AWSV4SignerAuth
            except ImportError as e:
                raise RuntimeError(
                    "OpenSearch SigV4 auth requires boto3 and opensearch-py. "
                    "Install with: knowledge-layer[opensearch]"
                ) from e

            credentials = boto3.Session(region_name=self.aws_region).get_credentials()
            if credentials is None:
                raise RuntimeError("No AWS credentials available for OpenSearch SigV4 auth")
            client_kwargs["http_auth"] = AWSV4SignerAuth(credentials, self.aws_region, self.aws_service)
        elif self.auth_type != "none":
            raise RuntimeError("OpenSearch auth_type must be one of: none, basic, sigv4")

        if self.auth_type == "sigv4" and parsed.hostname:
            host = {
                "host": parsed.hostname,
                "port": parsed.port or (443 if use_ssl else 80),
                "scheme": parsed.scheme or "https",
            }
            return OpenSearch(hosts=[host], **client_kwargs)

        return OpenSearch(hosts=[self.endpoint], **client_kwargs)

    def _get_client(self):
        with self._client_lock:
            if self._client is None:
                self._client = self._create_client()
            return self._client

    def _index_mapping(self, collection_name: str, description: str | None = None) -> dict[str, Any]:
        now = _utc_now().isoformat()
        meta = {
            "backend": "opensearch",
            "collection_name": collection_name,
            "description": description,
            "created_at": now,
            "updated_at": now,
            "embedding_model": self.embed_model_name,
            "embedding_dim": self.embedding_dim,
        }
        return {
            "settings": {
                "index": {
                    "knn": True,
                    "knn.algo_param.ef_search": self.ef_search,
                }
            },
            "mappings": {
                "_meta": meta,
                "dynamic_templates": [
                    {
                        "metadata_strings": {
                            "path_match": "metadata.*",
                            "match_mapping_type": "string",
                            "mapping": {"type": "keyword", "ignore_above": 1024},
                        }
                    }
                ],
                "properties": {
                    "chunk_id": {"type": "keyword"},
                    "file_id": {"type": "keyword"},
                    "file_name": {"type": "keyword"},
                    self.text_field: {"type": "text"},
                    self.vector_field: {
                        "type": "knn_vector",
                        "dimension": self.embedding_dim,
                        "method": {
                            "name": "hnsw",
                            "space_type": self.space_type,
                            "engine": self.engine,
                            "parameters": {
                                "ef_construction": self.ef_construction,
                                "m": self.m,
                            },
                        },
                    },
                    "display_citation": {"type": "keyword"},
                    "page_number": {"type": "integer"},
                    "content_type": {"type": "keyword"},
                    "content_subtype": {"type": "keyword"},
                    "file_size": {"type": "long"},
                    "metadata": {"type": "object", "enabled": True},
                    "created_at": {"type": "date"},
                    "updated_at": {"type": "date"},
                },
            },
        }

    def _get_index_meta(self, index_name: str) -> dict[str, Any]:
        client = self._get_client()
        try:
            info = client.indices.get(index=index_name)
            index_info = info.get(index_name, {}) if isinstance(info, dict) else {}
            return (index_info.get("mappings") or {}).get("_meta") or {}
        except Exception:
            return {}

    def _put_index_meta(self, index_name: str, meta: dict[str, Any]) -> None:
        client = self._get_client()
        try:
            client.indices.put_mapping(index=index_name, body={"_meta": meta})
        except Exception as e:
            logger.debug("Failed to update OpenSearch mapping metadata for %s: %s", index_name, e)

    def _ensure_index(self, collection_name: str, description: str | None = None) -> str:
        client = self._get_client()
        index_name = self._index_name_for_collection(collection_name)
        if not client.indices.exists(index=index_name):
            client.indices.create(index=index_name, body=self._index_mapping(collection_name, description))
        return index_name

    def _update_collection_timestamp(self, collection_name: str) -> None:
        index_name = self._index_name_for_collection(collection_name)
        meta = self._get_index_meta(index_name)
        if not meta:
            return
        meta["updated_at"] = _utc_now().isoformat()
        self._put_index_meta(index_name, meta)

    def _health_check_client(self) -> bool:
        client = self._get_client()
        if client.ping():
            return True

        if self.auth_type == "sigv4" and self.aws_service == "aoss":
            client.transport.perform_request("GET", "/_cat/indices")
            return True

        return False


@register_ingestor("opensearch")
class OpenSearchIngestor(TTLCleanupMixin, _OpenSearchConfigMixin, BaseIngestor):
    """OpenSearch-backed document ingestor."""

    backend_name = "opensearch"

    def __init__(self, config: dict[str, Any] | None = None):
        super().__init__(config)
        self._configure_opensearch()

        self.embed_model_name = self.config.get("embed_model", DEFAULT_EMBED_MODEL)
        self.embed_base_url = self.config.get("embed_base_url", DEFAULT_EMBED_BASE_URL)
        self.embedding_batch_size = int(self.config.get("embedding_batch_size", DEFAULT_EMBEDDING_BATCH_SIZE))
        self.chunk_size = int(self.config.get("chunk_size", DEFAULT_CHUNK_SIZE))
        self.chunk_overlap = int(self.config.get("chunk_overlap", DEFAULT_CHUNK_OVERLAP))
        self.generate_summary_enabled = self.config.get("generate_summary", False)
        self.summary_llm = self.config.get("summary_llm")
        self.ingestion_mode = str(self.config.get("ingestion_mode", DEFAULT_INGESTION_MODE)).lower()
        self.dask_scheduler_address = self.config.get("dask_scheduler_address", DEFAULT_DASK_SCHEDULER_ADDRESS)
        self.dask_file_transfer = str(self.config.get("dask_file_transfer", DEFAULT_DASK_FILE_TRANSFER)).lower()

        self._jobs: dict[str, IngestionJobStatus] = {}
        self._files: dict[str, FileInfo] = {}
        self._lock = threading.RLock()

        if self.config.get("start_ttl_cleanup", True):
            self._start_ttl_cleanup_task(COLLECTION_TTL_HOURS, TTL_CLEANUP_INTERVAL_SECONDS)
        logger.info("OpenSearchIngestor initialized: endpoint=%s, auth_type=%s", self.endpoint, self.auth_type)

    def _embed_texts(self, texts: list[str]) -> list[list[float]]:
        try:
            from openai import OpenAI
        except ImportError as e:
            raise RuntimeError(
                "OpenSearch ingestion requires openai for embeddings. Install knowledge-layer[opensearch]."
            ) from e

        client = OpenAI(base_url=self.embed_base_url, api_key=_resolve_embedding_api_key(self.embed_base_url))
        embeddings: list[list[float]] = []
        for start in range(0, len(texts), self.embedding_batch_size):
            batch = texts[start : start + self.embedding_batch_size]
            response = client.embeddings.create(
                model=self.embed_model_name,
                input=batch,
                extra_body={"input_type": "passage"},
            )
            embeddings.extend([list(item.embedding) for item in response.data])
        return embeddings

    def submit_job(
        self,
        file_paths: list[str],
        collection_name: str,
        config: dict[str, Any] | None = None,
    ) -> str:
        """Submit an ingestion job and return immediately with a polling job ID."""
        job_id = str(uuid.uuid4())
        job_config = {**self.config, **(config or {})}
        original_filenames = job_config.get("original_filenames", [])
        requested_file_id = job_config.get("file_id")
        file_metadata = job_config.get("metadata") or {}

        validated_paths = [path for path in file_paths if os.path.exists(path)]
        if not validated_paths:
            job = IngestionJobStatus(
                job_id=job_id,
                status=JobState.FAILED,
                submitted_at=_utc_now(),
                completed_at=_utc_now(),
                total_files=len(file_paths),
                processed_files=0,
                collection_name=collection_name,
                backend=self.backend_name,
                error_message="No valid file paths provided",
            )
            with self._lock:
                self._jobs[job_id] = job
            return job_id

        file_details = []
        for i, path in enumerate(validated_paths):
            file_name = original_filenames[i] if i < len(original_filenames) else Path(path).name
            file_id = requested_file_id if requested_file_id and len(validated_paths) == 1 else str(uuid.uuid4())
            file_details.append(
                FileProgress(
                    file_id=file_id,
                    file_name=file_name,
                    status=FileStatus.UPLOADING,
                    progress_percent=0.0,
                )
            )
            with self._lock:
                self._files[file_id] = FileInfo(
                    file_id=file_id,
                    file_name=file_name,
                    collection_name=collection_name,
                    status=FileStatus.UPLOADING,
                    file_size=os.path.getsize(path),
                    uploaded_at=_utc_now(),
                    metadata={**file_metadata, "job_id": job_id},
                )

        job = IngestionJobStatus(
            job_id=job_id,
            status=JobState.PENDING,
            submitted_at=_utc_now(),
            total_files=len(validated_paths),
            processed_files=0,
            collection_name=collection_name,
            backend=self.backend_name,
            file_details=file_details,
        )
        with self._lock:
            self._jobs[job_id] = job

        if self._should_use_dask_ingestion():
            try:
                self._start_dask_ingestion(job_id, validated_paths, collection_name, job_config)
            except Exception as e:
                if self.ingestion_mode == "auto":
                    logger.warning("Falling back to local OpenSearch ingestion because Dask submit failed: %s", e)
                    self._start_local_ingestion(job_id, validated_paths, collection_name, job_config)
                else:
                    self._mark_job_failed(job_id, f"Dask ingestion submission failed: {e}")
                    if job_config.get("cleanup_files", False):
                        self._cleanup_paths(validated_paths)
        else:
            self._start_local_ingestion(job_id, validated_paths, collection_name, job_config)
        return job_id

    def _should_use_dask_ingestion(self) -> bool:
        if self.ingestion_mode == "local":
            return False
        if self.ingestion_mode == "dask":
            return True
        if self.ingestion_mode == "auto":
            return bool(self.dask_scheduler_address)
        raise RuntimeError("OpenSearch ingestion_mode must be one of: local, dask, auto")

    def _start_local_ingestion(
        self,
        job_id: str,
        file_paths: list[str],
        collection_name: str,
        job_config: dict[str, Any],
    ) -> None:
        thread = threading.Thread(
            target=self._run_ingestion,
            args=(job_id, file_paths, collection_name, job_config),
            daemon=True,
        )
        thread.start()

    def _create_dask_client(self):
        if not self.dask_scheduler_address:
            raise RuntimeError(
                "Dask ingestion requires OPENSEARCH_DASK_SCHEDULER_ADDRESS or NAT_DASK_SCHEDULER_ADDRESS"
            )
        try:
            from distributed import Client
        except ImportError as e:
            raise RuntimeError("Dask ingestion requires the distributed package") from e
        return Client(self.dask_scheduler_address, timeout=f"{self.timeout}s")

    def _start_dask_ingestion(
        self,
        job_id: str,
        file_paths: list[str],
        collection_name: str,
        job_config: dict[str, Any],
    ) -> None:
        from knowledge_layer.opensearch.distributed import run_opensearch_ingestion_task

        with self._lock:
            job = self._jobs[job_id]
            job.status = JobState.PROCESSING
            job.started_at = _utc_now()
            job.metadata["ingestion_mode"] = "dask"
            for detail in job.file_details:
                detail.status = FileStatus.INGESTING
                tracked = self._files.get(detail.file_id)
                if tracked:
                    tracked.status = FileStatus.INGESTING

        payloads = self._build_dask_file_payloads(job_id, file_paths, job_config)
        worker_config = self._worker_config(job_config)
        client = self._create_dask_client()
        future = client.submit(
            run_opensearch_ingestion_task,
            worker_config,
            payloads,
            collection_name,
            key=f"aiq-opensearch-ingest-{job_id}",
            pure=False,
        )
        thread = threading.Thread(
            target=self._monitor_dask_ingestion,
            args=(job_id, future, client, file_paths, job_config),
            daemon=True,
        )
        thread.start()

    def _worker_config(self, job_config: dict[str, Any]) -> dict[str, Any]:
        worker_config = dict(job_config)
        worker_config.pop("summary_llm", None)
        worker_config["start_ttl_cleanup"] = False
        worker_config["generate_summary"] = False
        return worker_config

    def _build_dask_file_payloads(
        self,
        job_id: str,
        file_paths: list[str],
        job_config: dict[str, Any],
    ) -> list[dict[str, Any]]:
        with self._lock:
            job = self._jobs[job_id]
            details = list(job.file_details)
        file_metadata = job_config.get("metadata") or {}
        payloads = []
        for path, detail in zip(file_paths, details, strict=True):
            payload = {
                "file_id": detail.file_id,
                "file_name": detail.file_name,
                "metadata": file_metadata,
            }
            if self.dask_file_transfer == "bytes":
                payload["data"] = Path(path).read_bytes()
                payload["suffix"] = Path(path).suffix
            elif self.dask_file_transfer == "paths":
                payload["path"] = path
            else:
                raise RuntimeError("OpenSearch dask_file_transfer must be one of: bytes, paths")
            payloads.append(payload)
        return payloads

    def _monitor_dask_ingestion(
        self,
        job_id: str,
        future: Any,
        client: Any,
        file_paths: list[str],
        job_config: dict[str, Any],
    ) -> None:
        try:
            result = future.result()
            self._apply_dask_ingestion_result(job_id, result)
        except Exception as e:
            logger.exception("OpenSearch Dask ingestion job failed")
            self._mark_job_failed(job_id, str(e))
        finally:
            close = getattr(client, "close", None)
            if close:
                close()
            if job_config.get("cleanup_files", False):
                self._cleanup_paths(file_paths)

    def _apply_dask_ingestion_result(self, job_id: str, result: dict[str, Any]) -> None:
        file_results = {item.get("file_id"): item for item in result.get("files", [])}
        with self._lock:
            job = self._jobs[job_id]
            for index, detail in enumerate(job.file_details):
                item = file_results.get(detail.file_id, {})
                status = FileStatus(item.get("status", FileStatus.FAILED))
                self._mark_file(
                    job,
                    index,
                    status,
                    chunks_created=int(item.get("chunks_created", 0)),
                    error=item.get("error_message"),
                )
                summary = item.get("summary")
                if summary:
                    from aiq_agent.knowledge import register_summary

                    register_summary(job.collection_name, detail.file_name, summary)
                    tracked = self._files.get(detail.file_id)
                    if tracked:
                        tracked.metadata["summary"] = summary

            failed_count = sum(1 for detail in job.file_details if detail.status == FileStatus.FAILED)
            job.processed_files = job.total_files
            job.completed_at = _utc_now()
            job.metadata.update(
                {
                    "index_name": result.get("index_name"),
                    "total_chunks": result.get("total_chunks", 0),
                    "embedding_model": result.get("embedding_model", self.embed_model_name),
                    "ingestion_mode": "dask",
                }
            )
            if failed_count == job.total_files:
                job.status = JobState.FAILED
                job.error_message = result.get("error_message") or "All files failed ingestion"
            else:
                job.status = JobState.COMPLETED
                job.error_message = None
        self._update_collection_timestamp(job.collection_name)

    def _mark_job_failed(self, job_id: str, error: str) -> None:
        with self._lock:
            job = self._jobs[job_id]
            job.status = JobState.FAILED
            job.completed_at = _utc_now()
            job.error_message = error
            job.processed_files = job.total_files
            for index, _ in enumerate(job.file_details):
                self._mark_file(job, index, FileStatus.FAILED, error=error)

    def _cleanup_paths(self, file_paths: list[str]) -> None:
        for file_path in file_paths:
            try:
                os.unlink(file_path)
            except OSError:
                pass

    def get_job_status(self, job_id: str) -> IngestionJobStatus:
        with self._lock:
            job = self._jobs.get(job_id)
            if job is None:
                return IngestionJobStatus(
                    job_id=job_id,
                    status=JobState.FAILED,
                    submitted_at=_utc_now(),
                    completed_at=_utc_now(),
                    total_files=0,
                    processed_files=0,
                    collection_name="unknown",
                    backend=self.backend_name,
                    error_message="Job ID not found",
                )
            return job.model_copy(deep=True)

    def create_collection(
        self,
        name: str,
        description: str | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> CollectionInfo:
        index_name = self._ensure_index(name, description)
        meta = self._get_index_meta(index_name)
        if metadata:
            meta.update(metadata)
            self._put_index_meta(index_name, meta)
        return self._collection_info_from_index(name, index_name, meta)

    def delete_collection(self, name: str) -> bool:
        client = self._get_client()
        index_name = self._index_name_for_collection(name)
        try:
            if not client.indices.exists(index=index_name):
                return False
            client.indices.delete(index=index_name)
            from aiq_agent.knowledge import clear_collection_summaries

            clear_collection_summaries(name)
            with self._lock:
                self._files = {fid: f for fid, f in self._files.items() if f.collection_name != name}
            return True
        except Exception as e:
            logger.error("Failed to delete OpenSearch collection %s: %s", name, e)
            return False

    def list_collections(self) -> list[CollectionInfo]:
        client = self._get_client()
        pattern = f"{self.index_prefix}-*"
        try:
            response = client.indices.get(index=pattern)
        except Exception as e:
            logger.debug("Failed to list OpenSearch collections with pattern %s: %s", pattern, e)
            return []

        collections = []
        for index_name, info in response.items():
            meta = (info.get("mappings") or {}).get("_meta") or {}
            if meta.get("backend") != "opensearch":
                continue
            collection_name = meta.get("collection_name") or index_name.removeprefix(f"{self.index_prefix}-")
            collections.append(self._collection_info_from_index(collection_name, index_name, meta))
        return collections

    def get_collection(self, name: str) -> CollectionInfo | None:
        client = self._get_client()
        index_name = self._index_name_for_collection(name)
        try:
            if not client.indices.exists(index=index_name):
                return None
            return self._collection_info_from_index(name, index_name, self._get_index_meta(index_name))
        except Exception as e:
            logger.error("Failed to get OpenSearch collection %s: %s", name, e)
            return None

    def upload_file(
        self,
        file_path: str,
        collection_name: str,
        metadata: dict[str, Any] | None = None,
    ) -> FileInfo:
        path = Path(file_path)
        file_id = str(uuid.uuid4())
        if not path.exists():
            return FileInfo(
                file_id=file_id,
                file_name=path.name,
                collection_name=collection_name,
                status=FileStatus.FAILED,
                error_message=f"File not found: {file_path}",
            )

        job_id = self.submit_job(
            [file_path],
            collection_name,
            config={
                "file_id": file_id,
                "original_filenames": [path.name],
                "metadata": metadata or {},
            },
        )
        with self._lock:
            info = self._files[file_id]
            info.status = FileStatus.INGESTING
            info.metadata["job_id"] = job_id
            return info.model_copy(deep=True)

    def delete_file(self, file_id: str, collection_name: str) -> bool:
        client = self._get_client()
        index_name = self._index_name_for_collection(collection_name)
        try:
            if not client.indices.exists(index=index_name):
                return False

            resolved_name = self._resolve_file_name(file_id, collection_name)
            body = {
                "query": {
                    "bool": {
                        "should": [
                            {"term": {"file_id": file_id}},
                            {"term": {"file_name": resolved_name}},
                            {"term": {"file_name": file_id}},
                        ],
                        "minimum_should_match": 1,
                    }
                }
            }
            if self.auth_type == "sigv4" and self.aws_service == "aoss":
                deleted = self._delete_file_documents_for_aoss(index_name, body)
            else:
                result = client.delete_by_query(index=index_name, body=body, refresh=True, conflicts="proceed")
                deleted = int(result.get("deleted", 0)) if isinstance(result, dict) else 0

            with self._lock:
                for tracked_id, tracked_file in list(self._files.items()):
                    if (
                        tracked_file.collection_name == collection_name
                        and tracked_file.file_name in (resolved_name, file_id)
                    ):
                        self._files.pop(tracked_id, None)
                    elif tracked_id == file_id:
                        self._files.pop(tracked_id, None)

            if deleted > 0:
                from aiq_agent.knowledge import unregister_summary

                unregister_summary(collection_name, resolved_name)
                self._update_collection_timestamp(collection_name)
                return True
            return False
        except Exception as e:
            logger.error("Failed to delete OpenSearch file %s: %s", file_id, e)
            return False

    def _delete_file_documents_for_aoss(self, index_name: str, query_body: dict[str, Any]) -> int:
        """AOSS doesn't support _delete_by_query; search first, then bulk delete by generated IDs."""
        client = self._get_client()
        deleted = 0

        seen_ids: set[str] = set()
        stale_batches = 0

        for _ in range(self.aoss_delete_max_batches):
            response = client.search(
                index=index_name,
                body={
                    "size": self.bulk_batch_size,
                    "_source": False,
                    "query": query_body["query"],
                },
                request_timeout=self.timeout,
            )
            hits = response.get("hits", {}).get("hits", [])
            if not hits:
                return deleted

            body = []
            for hit in hits:
                hit_id = hit.get("_id")
                if hit_id and hit_id not in seen_ids:
                    seen_ids.add(hit_id)
                    body.append({"delete": {"_index": index_name, "_id": hit_id}})

            if not body:
                stale_batches += 1
                if stale_batches >= 2:
                    return deleted
                time.sleep(self.aoss_delete_backoff_seconds)
                continue
            stale_batches = 0

            result = client.bulk(body=body, refresh=self.bulk_refresh, request_timeout=self.timeout)
            if isinstance(result, dict) and result.get("errors"):
                raise RuntimeError(f"OpenSearch bulk deletion failed: {result}")

            deleted += len(body)
            if self.aoss_delete_backoff_seconds:
                time.sleep(self.aoss_delete_backoff_seconds)

        raise RuntimeError(
            f"OpenSearch AOSS file deletion exceeded {self.aoss_delete_max_batches} search/delete batches"
        )

    def list_files(self, collection_name: str) -> list[FileInfo]:
        client = self._get_client()
        index_name = self._index_name_for_collection(collection_name)
        try:
            if not client.indices.exists(index=index_name):
                return []
            response = client.search(
                index=index_name,
                body={
                    "size": 10000,
                    "_source": [
                        "file_id",
                        "file_name",
                        "file_size",
                        "content_type",
                        "created_at",
                        "updated_at",
                        "metadata",
                    ],
                    "query": {"match_all": {}},
                },
            )
        except Exception as e:
            logger.error("Failed to list OpenSearch files for %s: %s", collection_name, e)
            return []

        files = self._files_from_hits(response, collection_name)
        existing_names = {f.file_name for f in files}
        with self._lock:
            for tracked in self._files.values():
                if (
                    tracked.collection_name == collection_name
                    and tracked.status == FileStatus.FAILED
                    and tracked.file_name not in existing_names
                ):
                    files.append(tracked.model_copy(deep=True))
                    existing_names.add(tracked.file_name)
        return files

    def get_file_status(self, file_id: str, collection_name: str) -> FileInfo | None:
        with self._lock:
            tracked = self._files.get(file_id)
            if tracked:
                job_id = tracked.metadata.get("job_id")
                if tracked.status == FileStatus.INGESTING and job_id:
                    job = self._jobs.get(job_id)
                    if job and job.status == JobState.COMPLETED:
                        tracked.status = FileStatus.SUCCESS
                        tracked.ingested_at = job.completed_at
                    elif job and job.status == JobState.FAILED:
                        tracked.status = FileStatus.FAILED
                        tracked.error_message = job.error_message
                return tracked.model_copy(deep=True)

        for file_info in self.list_files(collection_name):
            if file_id in (file_info.file_id, file_info.file_name):
                return file_info
        return None

    def generate_summary(self, text_content: str, file_name: str) -> str | None:
        if not self.generate_summary_enabled:
            return None
        return _generate_document_summary(text_content, file_name, self.summary_llm)

    async def health_check(self) -> bool:
        try:
            return self._health_check_client()
        except Exception as e:
            logger.warning("OpenSearch health check failed: %s", e)
            return False

    def _run_ingestion(
        self,
        job_id: str,
        file_paths: list[str],
        collection_name: str,
        config: dict[str, Any],
    ) -> None:
        try:
            with self._lock:
                job = self._jobs[job_id]
                job.status = JobState.PROCESSING
                job.started_at = _utc_now()
                for detail in job.file_details:
                    detail.status = FileStatus.INGESTING

            index_name = self._ensure_index(collection_name)
            original_filenames = config.get("original_filenames", [])
            total_chunks = 0

            for i, file_path in enumerate(file_paths):
                file_name = original_filenames[i] if i < len(original_filenames) else Path(file_path).name
                file_id = job.file_details[i].file_id
                try:
                    documents, summary_text = self._documents_for_file(
                        file_path,
                        file_id,
                        file_name,
                        config.get("metadata") or {},
                    )
                    if not documents:
                        self._mark_file(job, i, FileStatus.FAILED, error="No content extracted")
                        continue

                    embeddings = self._embed_texts([doc[self.text_field] for doc in documents])
                    for doc, embedding in zip(documents, embeddings, strict=True):
                        doc[self.vector_field] = embedding

                    self._bulk_index_documents(index_name, documents)
                    chunks_created = len(documents)
                    total_chunks += chunks_created
                    self._mark_file(job, i, FileStatus.SUCCESS, chunks_created=chunks_created)

                    if self.generate_summary_enabled:
                        summary = self.generate_summary(summary_text, file_name)
                        if summary:
                            from aiq_agent.knowledge import register_summary

                            register_summary(collection_name, file_name, summary)
                            with self._lock:
                                if file_id in self._files:
                                    self._files[file_id].metadata["summary"] = summary

                except Exception as e:
                    logger.exception("OpenSearch ingestion failed for %s", file_path)
                    self._mark_file(job, i, FileStatus.FAILED, error=str(e))

            with self._lock:
                failed_count = sum(1 for detail in job.file_details if detail.status == FileStatus.FAILED)
                job.processed_files = job.total_files
                job.completed_at = _utc_now()
                job.metadata = {
                    "index_name": index_name,
                    "total_chunks": total_chunks,
                    "embedding_model": self.embed_model_name,
                }
                if failed_count == job.total_files:
                    job.status = JobState.FAILED
                    job.error_message = "All files failed ingestion"
                else:
                    job.status = JobState.COMPLETED

            self._update_collection_timestamp(collection_name)

        except Exception as e:
            logger.exception("OpenSearch ingestion job failed")
            with self._lock:
                job = self._jobs[job_id]
                job.status = JobState.FAILED
                job.completed_at = _utc_now()
                job.error_message = str(e)
        finally:
            if config.get("cleanup_files", False):
                for file_path in file_paths:
                    try:
                        os.unlink(file_path)
                    except OSError:
                        pass

    def _documents_for_file(
        self,
        file_path: str,
        file_id: str,
        file_name: str,
        file_metadata: dict[str, Any] | None = None,
    ) -> tuple[list[dict[str, Any]], str]:
        file_size = os.path.getsize(file_path)
        now = _utc_now().isoformat()
        documents = []
        summary_parts = []
        file_metadata = file_metadata or {}

        for segment_text, page_number, segment_metadata in _read_file_segments(file_path):
            chunks = _chunk_text(segment_text, self.chunk_size, self.chunk_overlap)
            if not summary_parts and segment_text.strip():
                summary_parts.append(segment_text.strip()[:SUMMARY_MAX_INPUT_CHARS])

            for chunk_index, content in enumerate(chunks):
                content_type = ContentType.TEXT
                chunk_id = str(uuid.uuid4())
                if page_number:
                    display_citation = f"{file_name}, p.{page_number}"
                else:
                    display_citation = file_name

                documents.append(
                    {
                        "chunk_id": chunk_id,
                        "file_id": file_id,
                        "file_name": file_name,
                        self.text_field: content,
                        "display_citation": display_citation,
                        "page_number": page_number,
                        "content_type": content_type.value,
                        "content_subtype": None,
                        "file_size": file_size,
                        "metadata": {
                            **file_metadata,
                            **segment_metadata,
                            "chunk_index": chunk_index,
                            "source_path": file_path,
                        },
                        "created_at": now,
                        "updated_at": now,
                    }
                )

        return documents, "\n".join(summary_parts)

    def _bulk_index_documents(self, index_name: str, documents: list[dict[str, Any]]) -> None:
        client = self._get_client()
        for start in range(0, len(documents), self.bulk_batch_size):
            batch = documents[start : start + self.bulk_batch_size]
            body = []
            for doc in batch:
                action = {"index": {"_index": index_name}}
                if self.allow_document_ids:
                    action["index"]["_id"] = doc["chunk_id"]
                body.append(action)
                body.append(doc)
            result = client.bulk(body=body, refresh=self.bulk_refresh, request_timeout=self.timeout)
            if isinstance(result, dict) and result.get("errors"):
                raise RuntimeError(f"OpenSearch bulk indexing failed: {result}")

    def _mark_file(
        self,
        job: IngestionJobStatus,
        file_index: int,
        status: FileStatus,
        chunks_created: int = 0,
        error: str | None = None,
    ) -> None:
        with self._lock:
            if file_index < len(job.file_details):
                detail = job.file_details[file_index]
                detail.status = status
                detail.progress_percent = 100.0
                detail.chunks_created = chunks_created
                detail.error_message = error
                tracked = self._files.get(detail.file_id)
                if tracked:
                    tracked.status = status
                    tracked.chunk_count = chunks_created
                    tracked.error_message = error
                    if status == FileStatus.SUCCESS:
                        tracked.ingested_at = _utc_now()
            job.processed_files = min(job.total_files, file_index + 1)

    def _collection_info_from_index(
        self,
        collection_name: str,
        index_name: str,
        meta: dict[str, Any],
    ) -> CollectionInfo:
        client = self._get_client()
        chunk_count = 0
        try:
            count_result = client.count(index=index_name)
            chunk_count = int(count_result.get("count", 0))
        except Exception:
            pass

        files = self.list_files(collection_name) if chunk_count else []
        return CollectionInfo(
            name=collection_name,
            description=meta.get("description"),
            file_count=len(files),
            chunk_count=chunk_count,
            created_at=_parse_timestamp(meta.get("created_at")),
            updated_at=_parse_timestamp(meta.get("updated_at")),
            backend=self.backend_name,
            metadata={
                "index_name": index_name,
                "endpoint": self.endpoint,
                "embedding_model": meta.get("embedding_model", self.embed_model_name),
                "embedding_dim": meta.get("embedding_dim", self.embedding_dim),
            },
        )

    def _files_from_hits(self, response: dict[str, Any], collection_name: str) -> list[FileInfo]:
        grouped: dict[str, dict[str, Any]] = {}
        for hit in (response.get("hits") or {}).get("hits", []):
            source = hit.get("_source") or {}
            file_name = source.get("file_name", "unknown")
            entry = grouped.setdefault(
                file_name,
                {
                    "file_id": source.get("file_id") or file_name,
                    "chunk_count": 0,
                    "file_size": source.get("file_size"),
                    "content_types": set(),
                    "uploaded_at": _parse_timestamp(source.get("created_at")),
                    "ingested_at": _parse_timestamp(source.get("updated_at")),
                    "metadata": source.get("metadata") or {},
                },
            )
            entry["chunk_count"] += 1
            if source.get("content_type"):
                entry["content_types"].add(source["content_type"])

        files = []
        for file_name, info in grouped.items():
            files.append(
                FileInfo(
                    file_id=info["file_id"],
                    file_name=file_name,
                    collection_name=collection_name,
                    status=FileStatus.SUCCESS,
                    file_size=info["file_size"],
                    chunk_count=info["chunk_count"],
                    uploaded_at=info["uploaded_at"],
                    ingested_at=info["ingested_at"],
                    metadata={
                        **info["metadata"],
                        "content_types": sorted(info["content_types"]),
                    },
                )
            )
        return files

    def _resolve_file_name(self, file_id: str, collection_name: str) -> str:
        with self._lock:
            tracked = self._files.get(file_id)
            if tracked and tracked.collection_name == collection_name:
                return tracked.file_name
            for info in self._files.values():
                if info.collection_name == collection_name and info.file_name == file_id:
                    return info.file_name
        return file_id


@register_retriever("opensearch")
class OpenSearchRetriever(_OpenSearchConfigMixin, BaseRetriever):
    """OpenSearch-backed document retriever."""

    backend_name = "opensearch"

    def __init__(self, config: dict[str, Any] | None = None):
        super().__init__(config)
        self._configure_opensearch()
        self.embed_model_name = self.config.get("embed_model", DEFAULT_EMBED_MODEL)
        self.embed_base_url = self.config.get("embed_base_url", DEFAULT_EMBED_BASE_URL)
        self.embedding_batch_size = int(self.config.get("embedding_batch_size", DEFAULT_EMBEDDING_BATCH_SIZE))
        self.default_top_k = int(self.config.get("top_k", 10))
        logger.info("OpenSearchRetriever initialized: endpoint=%s, auth_type=%s", self.endpoint, self.auth_type)

    def _embed_texts(self, texts: list[str]) -> list[list[float]]:
        try:
            from openai import OpenAI
        except ImportError as e:
            raise RuntimeError(
                "OpenSearch retrieval requires openai for embeddings. Install knowledge-layer[opensearch]."
            ) from e

        client = OpenAI(base_url=self.embed_base_url, api_key=_resolve_embedding_api_key(self.embed_base_url))
        response = client.embeddings.create(
            model=self.embed_model_name,
            input=texts,
            extra_body={"input_type": "query"},
        )
        return [list(item.embedding) for item in response.data]

    async def retrieve(
        self,
        query: str,
        collection_name: str,
        top_k: int = 10,
        filters: dict[str, Any] | None = None,
    ) -> RetrievalResult:
        try:
            client = self._get_client()
            index_name = self._index_name_for_collection(collection_name)
            if not client.indices.exists(index=index_name):
                return RetrievalResult(
                    chunks=[],
                    query=query,
                    backend=self.backend_name,
                    success=False,
                    error_message=f"Collection '{collection_name}' not found",
                )

            query_embedding = self._embed_texts([query])[0]
            body = self._build_search_body(query_embedding, top_k or self.default_top_k, filters)
            response = client.search(index=index_name, body=body, request_timeout=self.timeout)
            chunks = [
                chunk for chunk in (self.normalize(hit) for hit in response.get("hits", {}).get("hits", [])) if chunk
            ]

            return RetrievalResult(
                chunks=chunks,
                total_tokens=sum(len(chunk.content.split()) for chunk in chunks),
                query=query,
                backend=self.backend_name,
                success=True,
            )
        except Exception as e:
            logger.error("OpenSearch retrieval failed: %s", e)
            return RetrievalResult(
                chunks=[],
                query=query,
                backend=self.backend_name,
                success=False,
                error_message=f"Retrieval failed: {str(e)[:100]}",
            )

    def _build_search_body(
        self,
        query_embedding: list[float],
        top_k: int,
        filters: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        knn_body: dict[str, Any] = {
            "vector": query_embedding,
            "k": top_k,
        }
        filter_query = self._build_filter_query(filters)
        if filter_query:
            knn_body["filter"] = filter_query

        return {
            "size": top_k,
            "_source": {
                "excludes": [self.vector_field],
            },
            "query": {
                "knn": {
                    self.vector_field: knn_body,
                }
            },
        }

    def _build_filter_query(self, filters: dict[str, Any] | None) -> dict[str, Any] | None:
        if not filters:
            return None
        if "filter" in filters and isinstance(filters["filter"], dict):
            return filters["filter"]

        clauses = []
        for key, value in filters.items():
            if key in ("filter", "filter_expr"):
                continue
            field_name = key if key in {"file_name", "content_type", "file_id"} else f"metadata.{key}"
            clauses.append({"term": {field_name: value}})
        if not clauses:
            return None
        return {"bool": {"filter": clauses}}

    def normalize(self, raw_result: Any) -> Chunk | None:
        if not isinstance(raw_result, dict):
            return None

        source = raw_result.get("_source") or {}
        content = source.get(self.text_field, "")
        file_name = source.get("file_name", "unknown")
        page_number = source.get("page_number")
        content_type = self._content_type_from_source(source)
        display_citation = source.get("display_citation") or self._display_citation(file_name, page_number)

        return Chunk(
            chunk_id=source.get("chunk_id") or raw_result.get("_id") or str(uuid.uuid4()),
            content=content or "",
            score=_score_to_similarity(raw_result.get("_score", 0.0)),
            file_name=file_name,
            page_number=page_number,
            display_citation=display_citation,
            content_type=content_type,
            content_subtype=source.get("content_subtype"),
            structured_data=source.get("structured_data"),
            image_storage_uri=source.get("image_storage_uri"),
            image_url=source.get("image_url"),
            metadata={
                **(source.get("metadata") or {}),
                "file_id": source.get("file_id"),
                "index": raw_result.get("_index"),
            },
        )

    async def health_check(self) -> bool:
        try:
            return self._health_check_client()
        except Exception:
            return False

    def _content_type_from_source(self, source: dict[str, Any]) -> ContentType:
        raw_type = str(source.get("content_type", "text")).lower()
        if raw_type == ContentType.TABLE.value:
            return ContentType.TABLE
        if raw_type == ContentType.CHART.value:
            return ContentType.CHART
        if raw_type == ContentType.IMAGE.value:
            return ContentType.IMAGE
        return ContentType.TEXT

    def _display_citation(self, file_name: str, page_number: Any) -> str:
        if page_number:
            return f"{file_name}, p.{page_number}"
        return file_name
